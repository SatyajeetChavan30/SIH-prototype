"""
Build the SIH 2026 idea-submission deck for PS 26161 (JalRaksha).

    python tools/sih-presentation/build_deck.py

Reads : SIH2026-IDEA-Presentation-Format.pptx   (official template, unmodified)
Writes: JalRaksha_SIH2026_Idea.pptx

Constraints taken from the template's own IMPORTANT INSTRUCTIONS slide:
  * Maximum six slides INCLUDING the title slide, so slide 7 is deleted.
  * The idea-details pointers are reproduced verbatim as section labels. Only
    slide 2's title is replaced, with the idea name, which is what it is for.
  * Points, diagrams and pictures — not paragraphs.
  * The portal takes PDF only; see export_pdf.ps1.

Every figure on these slides is a live number from this repository: the
validation gates come from GET /validation, the Tehri figures from
data/exports/264734ee.../run_summary.json and population_at_risk.json, and the
screenshots from capture_dashboard.py driving the running dashboard. Nothing
here is illustrative.
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent.parent
TEMPLATE = ROOT / "SIH2026-IDEA-Presentation-Format.pptx"
OUTPUT = ROOT / "JalRaksha_SIH2026_Idea.pptx"
IMG = HERE / "assets" / "prepared"
LOGOS = HERE / "assets" / "logos"

EMU_PER_INCH = 914400

# --------------------------------------------------------------------------- #
# Palette — deep water navy dominates, hydro teal supports, hazard red is the
# single sharp accent. Chosen to sit with the template's own blue furniture.
# --------------------------------------------------------------------------- #
NAVY = RGBColor(0x0F, 0x2B, 0x46)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
HAZARD = RGBColor(0xB0, 0x2E, 0x2A)
INK = RGBColor(0x1F, 0x29, 0x33)
MUTED = RGBColor(0x5A, 0x66, 0x72)
TINT = RGBColor(0xEE, 0xF4, 0xF9)
EDGE = RGBColor(0xC9, 0xDC, 0xEC)
WARM = RGBColor(0xFD, 0xF3, 0xE7)
WARM_EDGE = RGBColor(0xEC, 0xD3, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BODY = "Calibri"
TITLE_FONT = "Times New Roman"

TEAM_NAME = "<TEAM NAME>"
TEAM_ID = "<Team ID from SIH portal>"

# --------------------------------------------------------------------------- #
# Inline markup + text helpers. The **bold** / *italic* segmenter is carried
# over from build_ppt.py, which proved it against this same template.
# --------------------------------------------------------------------------- #

_BOLD = re.compile(r"\*\*(.+?)\*\*")
_ITAL = re.compile(r"(?<!\*)\*([^*]+?)\*(?!\*)")


def _segments(text: str):
    out, pos = [], 0
    for m in _BOLD.finditer(text):
        if m.start() > pos:
            out.append((text[pos : m.start()], False))
        out.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], False))

    final = []
    for chunk, bold in out:
        p = 0
        for m in _ITAL.finditer(chunk):
            if m.start() > p:
                final.append((chunk[p : m.start()], bold, False))
            final.append((m.group(1), bold, True))
            p = m.end()
        if p < len(chunk):
            final.append((chunk[p:], bold, False))
    return final


def write(
    tf,
    lines,
    *,
    size=10.0,
    color=INK,
    bullet=None,
    space_after=3.5,
    line_spacing=0.95,
    align=PP_ALIGN.LEFT,
    font=BODY,
    bold=False,
):
    """Fill a text frame. `lines` is a list of markup strings, one per paragraph."""
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(space_after)
        para.line_spacing = line_spacing
        text = f"{bullet} {line}" if bullet else line
        for chunk, b, italic in _segments(text):
            run = para.add_run()
            run.text = chunk
            run.font.size = Pt(size)
            run.font.name = font
            run.font.bold = bold or b
            run.font.italic = italic
            run.font.color.rgb = color
    return tf


def box(
    slide, x, y, w, h, *, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.08, lw=0.9
):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = adj
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.09)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    return sp


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = 0
    return tb


def label(slide, x, y, w, text, *, size=11.5, color=NAVY):
    """A mandated template pointer, reproduced verbatim as a section heading."""
    lb = textbox(slide, x, y, w, 0.30, anchor=MSO_ANCHOR.MIDDLE)
    write(lb.text_frame, [text], size=size, color=color, space_after=0, bold=True)
    return y + 0.34


def caption(slide, x, y, w, text, *, size=7.8, color=MUTED):
    cb = textbox(slide, x, y, w, 0.50)
    write(cb.text_frame, [text], size=size, color=color, space_after=1.5, line_spacing=0.98)
    return cb


def picture(slide, path: Path, x, y, w):
    """Place an image by width; returns its rendered height in inches."""
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return pic, pic.height / EMU_PER_INCH


def stat(slide, x, y, w, h, value, unit, cap_lines, *, accent=NAVY, vsize=23):
    """A large number with its unit and a small caption beneath."""
    box(slide, x, y, w, h, fill=TINT, line=EDGE)
    vb = textbox(slide, x + 0.13, y + 0.09, w - 0.26, 0.42)
    para = vb.text_frame.paragraphs[0]
    para.space_after = Pt(0)
    para.line_spacing = 0.92
    r = para.add_run()
    r.text = value
    r.font.size, r.font.name, r.font.bold, r.font.color.rgb = Pt(vsize), BODY, True, accent
    if unit:
        r2 = para.add_run()
        r2.text = " " + unit
        r2.font.size, r2.font.name, r2.font.color.rgb = Pt(9.5), BODY, MUTED
    cb = textbox(slide, x + 0.13, y + 0.54, w - 0.26, h - 0.62)
    write(cb.text_frame, cap_lines, size=8.0, color=MUTED, space_after=0.5, line_spacing=0.96)


def table(slide, x, y, w, rows, col_w, *, head_fill=NAVY, size=8.6, row_h=0.245):
    """A compact table. `rows[0]` is the header row."""
    n_rows, n_cols = len(rows), len(rows[0])
    shp = slide.shapes.add_table(
        n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows)
    )
    tbl = shp.table
    tbl.first_row = True
    tbl.horz_banding = False
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(row_h)
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.015)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = head_fill if i == 0 else (WHITE if i % 2 else TINT)
            write(
                cell.text_frame,
                [val],
                size=size,
                color=WHITE if i == 0 else INK,
                bold=(i == 0),
                space_after=0,
                align=PP_ALIGN.RIGHT if (j and i) else PP_ALIGN.LEFT,
            )
    return shp


def arrow_down(slide, cx, y, h=0.20, *, color=NAVY):
    sp = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(cx - 0.09), Inches(y), Inches(0.18), Inches(h)
    )
    sp.rotation = 180
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False


def arrow_right(slide, x, cy, w=0.22, *, color=NAVY):
    sp = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(cy - 0.09), Inches(w), Inches(0.18)
    )
    sp.rotation = 90
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False


def arrow_up(slide, cx, y, h=0.18, *, color=NAVY):
    sp = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(cx - 0.09), Inches(y), Inches(0.18), Inches(h)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False


def arrow_left(slide, x, cy, w=0.22, *, color=NAVY):
    sp = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(cy - 0.09), Inches(w), Inches(0.18)
    )
    sp.rotation = 270
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False


def line_h(slide, x, y, w, *, color=NAVY, t=0.022):
    """Thin horizontal connector — the fork and the feedback loop are built from these."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(t))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False


def line_v(slide, x, y, h, *, color=NAVY, t=0.022):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(t), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False


def swimlane(slide, x, y, w, h, title, *, accent=NAVY, fill=TINT, edge=EDGE):
    """A grouped container with a title strip — the Fig 4.1 architecture idiom."""
    box(slide, x, y, w, h, fill=fill, line=edge)
    tb = textbox(slide, x + 0.10, y + 0.06, w - 0.20, 0.22, anchor=MSO_ANCHOR.MIDDLE)
    write(
        tb.text_frame,
        [title],
        size=8.2,
        color=accent,
        align=PP_ALIGN.CENTER,
        space_after=0,
        bold=True,
    )
    return y + 0.32


def unit(slide, x, y, w, h, head, sub, *, accent=NAVY, shape=MSO_SHAPE.RECTANGLE, dashed=False):
    """A component box: bold name over a parenthetical qualifier."""
    sp = box(slide, x, y, w, h, fill=WHITE, line=accent, shape=shape, lw=0.85)
    if dashed:
        sp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    inner = 0.16 if shape == MSO_SHAPE.HEXAGON else 0.06
    tb = textbox(slide, x + inner, y + 0.05, w - 2 * inner, h - 0.10, anchor=MSO_ANCHOR.MIDDLE)
    lines = [f"**{head}**"] + ([sub] if sub else [])
    write(
        tb.text_frame,
        lines,
        size=7.4,
        color=INK,
        align=PP_ALIGN.CENTER,
        space_after=0.6,
        line_spacing=0.94,
    )
    return sp


def edge(slide, x, y, w, text, *, color=NAVY):
    """A labelled connector between two architecture groups."""
    cy = y
    line_h(slide, x, cy, w - 0.13, color=color, t=0.018)
    arrow_right(slide, x + w - 0.15, cy + 0.009, w=0.15, color=color)
    lb = textbox(slide, x - 0.06, cy - 0.36, w + 0.12, 0.32)
    write(
        lb.text_frame,
        [text],
        size=6.3,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        space_after=0,
        line_spacing=0.92,
    )


def logo_grid(slide, x, y, w, marks, *, cols=4, cell_h=0.62, icon=0.30):
    """Technology-stack marks, generated by logos.js from simple-icons."""
    cw = w / cols
    for i, (slug, lab) in enumerate(marks):
        row, col = divmod(i, cols)
        cx, cy = x + col * cw, y + row * cell_h
        slide.shapes.add_picture(
            str(LOGOS / f"{slug}.png"),
            Inches(cx + (cw - icon) / 2),
            Inches(cy),
            height=Inches(icon),
        )
        lb = textbox(slide, cx, cy + icon + 0.05, cw, 0.16)
        write(lb.text_frame, [lab], size=6.4, color=MUTED, align=PP_ALIGN.CENTER, space_after=0)
    return y + cell_h * ((len(marks) + cols - 1) // cols)


def tech_group(slide, x, y, w, cat, items, *, accent=NAVY, h=0.50, chip_w=1.16, size=8.4):
    """A category chip with its technologies beneath. Returns the next y."""
    sq = box(slide, x, y, chip_w, 0.22, fill=accent, adj=0.22)
    tf = sq.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(tf, [cat], size=7.8, color=WHITE, align=PP_ALIGN.CENTER, space_after=0, bold=True)
    tb = textbox(slide, x, y + 0.25, w, h)
    write(tb.text_frame, [items], size=size, color=INK, space_after=0, line_spacing=0.97)
    return y + 0.25 + h


def flow_node(slide, x, y, w, h, num, head, body, *, accent=NAVY, fill=TINT, edge=EDGE):
    """The deck's motif: a numbered chip heading a tinted block."""
    box(slide, x, y, w, h, fill=fill, line=edge)
    sq = box(slide, x + 0.09, y + 0.09, 0.26, 0.26, fill=accent, adj=0.2)
    tf = sq.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(tf, [str(num)], size=9.5, color=WHITE, align=PP_ALIGN.CENTER, space_after=0, bold=True)
    hb = textbox(slide, x + 0.42, y + 0.10, w - 0.52, 0.24, anchor=MSO_ANCHOR.MIDDLE)
    write(hb.text_frame, [head], size=9.4, color=accent, space_after=0, bold=True)
    if body:
        bb = textbox(slide, x + 0.12, y + 0.38, w - 0.24, h - 0.46)
        write(bb.text_frame, body, size=8.1, color=INK, space_after=1.2, line_spacing=0.94)


# --------------------------------------------------------------------------- #
# Template surgery
# --------------------------------------------------------------------------- #


def drop_shape(slide, name):
    for sh in list(slide.shapes):
        if sh.name == name:
            sh._element.getparent().remove(sh._element)
            return True
    return False


def _reset(tf):
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    tf.paragraphs[0].clear()


def set_shape_text(slide, name, lines, **kw):
    for sh in slide.shapes:
        if sh.name == name and sh.has_text_frame:
            _reset(sh.text_frame)
            write(sh.text_frame, lines, **kw)
            return sh
    return None


def stamp_team(slide):
    for sh in slide.shapes:
        if sh.name.startswith("Oval") and sh.has_text_frame:
            _reset(sh.text_frame)
            write(
                sh.text_frame,
                [TEAM_NAME],
                size=8,
                color=INK,
                align=PP_ALIGN.CENTER,
                space_after=0,
                bold=True,
            )


def retitle(slide, text, size=30):
    for sh in slide.shapes:
        if sh.name.startswith("Title"):
            _reset(sh.text_frame)
            run = sh.text_frame.paragraphs[0].add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.name = TITLE_FONT
            return


def delete_slide(prs, index):
    """Remove a slide by zero-based index, dropping its relationship too."""
    sld_id_lst = prs.slides._sldIdLst
    entries = list(sld_id_lst)
    prs.part.drop_rel(entries[index].rId)
    sld_id_lst.remove(entries[index])


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #


def logo_mark(slide, x, y, d, *, ring=NAVY, water=TEAL):
    """The JalRaksha mark: a dam wall holding water, inside a navy disc."""
    disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    disc.fill.solid()
    disc.fill.fore_color.rgb = ring
    disc.line.fill.background()
    disc.shadow.inherit = False

    # Impounded water: a wave band across the lower half of the disc.
    wave = slide.shapes.add_shape(
        MSO_SHAPE.DOUBLE_WAVE,
        Inches(x + d * 0.13),
        Inches(y + d * 0.50),
        Inches(d * 0.74),
        Inches(d * 0.30),
    )
    wave.fill.solid()
    wave.fill.fore_color.rgb = water
    wave.line.fill.background()
    wave.shadow.inherit = False

    # The dam wall it is held behind.
    wall = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + d * 0.45),
        Inches(y + d * 0.20),
        Inches(d * 0.10),
        Inches(d * 0.46),
    )
    wall.fill.solid()
    wall.fill.fore_color.rgb = WHITE
    wall.line.fill.background()
    wall.shadow.inherit = False
    return disc


def wordmark(slide, x, y, w, *, size=32):
    """Two-tone JalRaksha wordmark: 'Jal' navy, 'Raksha' teal."""
    tb = textbox(slide, x, y, w, size / 60.0 + 0.16, anchor=MSO_ANCHOR.MIDDLE)
    para = tb.text_frame.paragraphs[0]
    para.space_after = Pt(0)
    para.line_spacing = 0.95
    for text, colour in (("Jal", NAVY), ("Raksha", TEAL)):
        run = para.add_run()
        run.text = text
        run.font.size, run.font.name, run.font.bold = Pt(size), BODY, True
        run.font.color.rgb = colour
    return tb


def swatches(slide, x, y, pairs, *, sq=0.22, gap=0.92):
    """Accent-colour chips with their hex, the way a brand sheet shows them."""
    for i, (colour, hexname) in enumerate(pairs):
        cx = x + i * gap
        box(slide, cx, y, sq, sq, fill=colour, adj=0.2)
        lb = textbox(slide, cx + sq + 0.05, y, gap - sq - 0.10, sq, anchor=MSO_ANCHOR.MIDDLE)
        write(lb.text_frame, [hexname], size=7.4, color=MUTED, space_after=0)


def slide1(s):
    # The template's plain subtitle is replaced by a proper brand lockup —
    # mark, two-tone wordmark and tagline — so the product reads as a product.
    drop_shape(s, "Subtitle 3")
    logo_mark(s, 0.40, 1.16, 0.82)
    wordmark(s, 1.34, 1.18, 4.60, size=32)
    tb = textbox(s, 1.38, 1.74, 5.60, 0.24)
    write(
        tb.text_frame,
        ["*Dam-break screening you can check.*"],
        size=10.5,
        color=TEAL,
        space_after=0,
    )

    fields = [
        ("Problem Statement ID", "26161"),
        (
            "Problem Statement Title",
            "Dam Break Inundation Modelling Using Hydrodynamic Modelling of any River",
        ),
        ("Theme", "Disaster Management"),
        ("PS Category", "Software"),
        ("Team ID", TEAM_ID),
        ("Team Name (Registered on portal)", TEAM_NAME),
    ]
    set_shape_text(
        s,
        "TextBox 9",
        [f"**{k} —**  {v}" for k, v in fields],
        size=13,
        color=INK,
        space_after=10,
        line_spacing=0.98,
    )
    lb = textbox(s, 0.40, 5.30, 3.00, 0.22)
    write(lb.text_frame, ["**THEME**"], size=8.0, color=NAVY, space_after=0)
    swatches(
        s,
        0.40,
        5.56,
        [(NAVY, "#0F2B46"), (TEAL, "#0E7C7B"), (HAZARD, "#B02E2A")],
    )

    tb = textbox(s, 0.36, 6.15, 6.40, 0.95)
    write(
        tb.text_frame,
        [
            "**Near-field SPH × far-field 2D shallow water, cross-checked against a real "
            "Delft3D FM kernel** — arrival-time maps for any Indian dam from nothing but "
            "lat/lon, height and storage. Open data only. No ground survey.",
        ],
        size=10.5,
        color=MUTED,
        line_spacing=1.0,
    )


def slide2(s):
    # The template's title placeholder spans x=0.20..12.20, which runs under the
    # team-name oval on the left and the SIH logo on the right. Re-seat it
    # between the two rather than letting either clip the idea title.
    retitle(s, "JALRAKSHA — dam-break screening, with the model checked", size=21)
    for sh in s.shapes:
        if sh.name.startswith("Title"):
            sh.left, sh.width = Inches(1.90), Inches(8.65)
    drop_shape(s, "TextBox 8")
    stamp_team(s)

    # ---- left: what it is -------------------------------------------------
    x, w = 0.35, 3.12
    y = label(s, x, 1.34, w, "Detailed explanation of the proposed solution")
    box(s, x, y, w, 5.08, fill=TINT, line=EDGE)
    tb = textbox(s, x + 0.12, y + 0.11, w - 0.24, 4.86)
    write(
        tb.text_frame,
        [
            "**Operator enters four things** — lat/lon, dam height, gross storage, breach "
            "mode. Nothing else.",
            "**Data fetches itself:** Copernicus GLO-30 DEM from public AWS COGs, GHSL "
            "population, Sentinel-1 SAR via Earth Engine, ESA WorldCover roughness.",
            "**Breach hydrograph** from a 4-regression Monte-Carlo ensemble — Costa 1985, "
            "Froehlich 1995, MacDonald–Langridge 1984, Von Thun–Gillette 1990.",
            "**Far field:** in-house 2D shallow-water finite volume — HLLC Riemann solver, "
            "Audusse hydrostatic reconstruction, MUSCL, wetting and drying.",
            "**Near field:** Smoothed Particle Hydrodynamics resolves the non-hydrostatic "
            "breach jet, where depth-averaging is invalid. One-way handoff, ~600 m over 15 s.",
            "**Outputs** depth, velocity and arrival-time rasters to a browser dashboard, "
            "plus GeoTIFF, .SHP and .KML/KMZ.",
        ],
        size=9.9,
        bullet="▪",
        space_after=8.5,
    )

    # ---- centre: the working thing ---------------------------------------
    cx, cw = 3.60, 5.78
    label(s, cx, 1.34, cw, "The prototype, running")
    _, ph = picture(s, IMG / "workspace.png", cx, 1.70, cw)
    caption(
        s,
        cx,
        1.70 + ph + 0.06,
        cw,
        "Live dashboard — Tehri (260 m, 3,540 MCM). Left: Leaflet inundation envelope with "
        "hazard classes over the Bhagirathi. Right: the same envelope on Cesium 3D over "
        "Copernicus terrain, with gauges labelled by downstream distance. 600×600 cells at "
        "200 m, 60 km radius, 3 h simulated, EPSG:32644.",
    )

    # ---- right: how it addresses the problem, and what is new -------------
    rx, rw = 9.52, 3.43
    y = label(s, rx, 1.34, rw, "How it addresses the problem")
    box(s, rx, y, rw, 2.28, fill=TINT, line=EDGE)
    tb = textbox(s, rx + 0.12, y + 0.10, rw - 0.24, 2.08)
    write(
        tb.text_frame,
        [
            "Engineered dam break **and** natural landslide-dam blockage, one pipeline.",
            "GUI for input and tiled output; **.shp / .kml export** exactly as the PS asks.",
            "Any DEM, any AOI, any Indian dam — five presets plus custom coordinates.",
            "Sentinel-1 SAR loop scores simulated extent against **observed** flood extent.",
            "Runs **offline from cache** after first fetch — demo-day networks assumed hostile.",
        ],
        size=9.0,
        bullet="▪",
        space_after=5,
    )

    y2 = label(s, rx, y + 2.40, rw, "Innovation and uniqueness of the solution")
    box(s, rx, y2, rw, 2.62, fill=WARM, line=WARM_EDGE)
    lead = textbox(s, rx + 0.12, y2 + 0.08, rw - 0.24, 0.20)
    write(
        lead.text_frame,
        ["Against how Tier-1 screening is done today:"],
        size=8.4,
        color=INK,
        space_after=0,
    )
    table(
        s,
        rx + 0.12,
        y2 + 0.30,
        rw - 0.24,
        [
            ["Current practice", "JalRaksha"],
            ["Weeks of specialist setup", "Minutes, 4 inputs"],
            ["Licence + field survey", "Free data, no survey"],
            ["No independent check", "Scored vs Delft3D FM"],
            ["One number for peak flow", "p05–p95 band"],
            ["Bad SAR scenes drawn", "Refused at a 0.50 gate"],
        ],
        [1.55, 1.64],
        head_fill=HAZARD,
        size=7.7,
        row_h=0.225,
    )
    foot = textbox(s, rx + 0.12, y2 + 1.72, rw - 0.24, 0.80)
    write(
        foot.text_frame,
        [
            "**The cross-check is the product.** The same Ritter dam-break runs through our "
            "solver and through a real Deltares kernel, and both are scored against the exact "
            "solution. Most entries have no independent check at all.",
        ],
        size=8.4,
        space_after=0,
        line_spacing=0.96,
    )


def slide3(s):
    drop_shape(s, "TextBox 8")
    stamp_team(s)

    # ======================================================================
    # Top band — technologies. Three groups across, then the stack as marks.
    # ======================================================================
    label(s, 0.35, 1.28, 12.60, "Technologies to be used")
    box(s, 0.35, 1.62, 12.60, 1.52, fill=TINT, line=EDGE)

    gw = 4.05
    for i, (cat, accent, items) in enumerate(
        [
            (
                "SOLVERS",
                NAVY,
                "In-house **2D shallow-water finite volume** \u2014 HLLC flux \u00b7 Audusse hydrostatic "
                "reconstruction \u00b7 MUSCL/minmod on \u03b7 = b + h. **SPH** near field. **Delft3D FM** "
                "(dflowfm-cli) as reference kernel.",
            ),
            (
                "OPEN DATA",
                TEAL,
                "Copernicus **GLO-30** from public AWS COGs \u00b7 **GHSL P2023A** \u00b7 **Sentinel-1 GRD** "
                "\u00b7 JRC Global Surface Water \u00b7 ESA WorldCover \u00b7 CWC NRLD/NRSD. None login-gated.",
            ),
            (
                "PROOF",
                HAZARD,
                "Monte-Carlo over four breach regressions with Wahl bands. **pytest gates** \u2014 "
                "analytical, lake-at-rest, mass conservation and the Delft3D FM cross-check, all "
                "currently passing.",
            ),
        ]
    ):
        tech_group(s, 0.50 + i * (gw + 0.15), 1.72, gw, cat, items, accent=accent, h=0.52, size=8.0)

    logo_grid(
        s,
        0.50,
        2.56,
        12.30,
        [
            ("python", "Python"),
            ("numpy", "NumPy"),
            ("numba", "Numba"),
            ("gdal", "GDAL"),
            ("fastapi", "FastAPI"),
            ("sqlite", "SQLite"),
            ("react", "React"),
            ("leaflet", "Leaflet"),
            ("cesium", "Cesium"),
            ("googleearthengine", "Earth Engine"),
            ("pytest", "pytest"),
        ],
        cols=11,
        cell_h=0.52,
        icon=0.26,
    )

    # ======================================================================
    # Bottom band — the architecture, read left to right.
    # ======================================================================
    label(s, 0.35, 3.24, 12.60, "Methodology and process for implementation")

    top, ch = 3.62, 2.82
    ax, aw = 0.35, 3.15
    hx, hw_ = 4.25, 1.45
    bx, bw = 6.45, 3.05
    cx, cw = 10.25, 2.70

    # ---- INGEST ----------------------------------------------------------
    iy = swimlane(s, ax, top, aw, ch, "INGEST  \u2014  open data, no login", accent=TEAL)
    uh, ug = 0.55, 0.06
    for j, (head, sub) in enumerate(
        [
            ("Operator input", "lat/lon \u00b7 height \u00b7 storage \u00b7 breach mode"),
            ("Copernicus GLO-30", "30 m DEM, public AWS COGs"),
            ("Google Earth Engine", "Sentinel-1 \u00b7 GHSL P2023A \u00b7 WorldCover"),
        ]
    ):
        unit(s, ax + 0.14, iy + j * (uh + ug), aw - 0.28, uh, head, sub, accent=TEAL)
    unit(
        s,
        ax + 0.60,
        iy + 3 * (uh + ug),
        aw - 1.20,
        0.50,
        "Offline cache",
        "COG + SQLite \u2014 runs from disk",
        accent=MUTED,
        shape=MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,
    )

    edge(s, ax + aw + 0.06, top + 1.30, hx - (ax + aw) - 0.12, "conditioned DEM\n+ Manning's n")

    # ---- the transform everything downstream depends on -------------------
    unit(
        s,
        hx,
        top + 0.72,
        hw_,
        1.30,
        "BREACH HYDROGRAPH",
        "Monte-Carlo over 4 published\nregressions \u2014 Q(t) with a\np05\u2013p95 band, never one number",
        accent=HAZARD,
        shape=MSO_SHAPE.HEXAGON,
        dashed=True,
    )

    edge(s, hx + hw_ + 0.06, top + 1.30, bx - (hx + hw_) - 0.12, "Q(t) at the breach\n229,952 m\u00b3/s median")

    # ---- SOLVE -----------------------------------------------------------
    sy = swimlane(s, bx, top, bw, ch, "SOLVE  \u2014  near field + far field", accent=NAVY)
    sh_, sg = 0.72, 0.08
    for j, (head, sub, acc) in enumerate(
        [
            ("Far field \u2014 2D SWE", "HLLC \u00b7 Audusse \u00b7 MUSCL\n600\u00d7600 cells at 200 m, 3 h", NAVY),
            ("Near field \u2014 SPH", "Lagrangian, non-hydrostatic\n~600 m over 15 s, one-way", TEAL),
            ("Delft3D FM", "dflowfm-cli \u2014 0.0349 m RMSE\nagainst our 0.0317 m", HAZARD),
        ]
    ):
        unit(s, bx + 0.14, sy + j * (sh_ + sg), bw - 0.28, sh_, head, sub, accent=acc)

    edge(s, bx + bw + 0.06, top + 1.30, cx - (bx + bw) - 0.12, "depth \u00b7 velocity\n\u00b7 arrival time")

    # ---- DELIVER ---------------------------------------------------------
    dy = swimlane(s, cx, top, cw, ch, "DELIVER  \u2014  decision products", accent=TEAL)
    for j, (head, sub) in enumerate(
        [
            ("Impact", "GHSL population at risk,\nbucketed by warning lead time"),
            ("Export", "Cloud-Optimized GeoTIFF\n\u00b7 .SHP \u00b7 .KML / KMZ"),
            ("Dashboard", "Leaflet 2D + Cesium 3D,\nrun picker works offline"),
        ]
    ):
        unit(s, cx + 0.14, dy + j * (sh_ + sg), cw - 0.28, sh_, head, sub, accent=TEAL)

    # ---- the validation loop, drawn as a return path ---------------------
    ly = top + ch + 0.16
    line_v(s, cx + cw / 2, top + ch, ly - (top + ch), color=TEAL)
    line_h(s, bx + bw / 2, ly, (cx + cw / 2) - (bx + bw / 2), color=TEAL)
    line_v(s, bx + bw / 2, top + ch, ly - (top + ch), color=TEAL)
    arrow_up(s, bx + bw / 2, top + ch - 0.13, h=0.14, color=TEAL)
    # The label sits to the left of the return path, pointing into it.
    lb = textbox(s, 0.40, ly - 0.11, (bx + bw / 2) - 0.55, 0.22, anchor=MSO_ANCHOR.MIDDLE)
    write(
        lb.text_frame,
        [
            "GEE VALIDATION LOOP \u21bb  Sentinel-1 observed extent scored against the simulation; "
            "scenes below 0.50 precision against JRC water are refused, not drawn.",
        ],
        size=6.8,
        color=TEAL,
        align=PP_ALIGN.RIGHT,
        space_after=0,
    )



def slide4(s):
    drop_shape(s, "TextBox 8")
    stamp_team(s)

    # ---- top left: feasibility, as evidence ------------------------------
    x, w = 0.35, 6.92
    y = label(s, x, 1.32, w, "Analysis of the feasibility of the idea")
    box(s, x, y, w, 2.56, fill=TINT, line=EDGE)

    tb = textbox(s, x + 0.12, y + 0.08, w - 0.24, 0.44)
    write(
        tb.text_frame,
        [
            "Not a proposal — a **built system**. These gates run against the live build, and "
            "are the same checks that block a merge in CI:",
        ],
        size=9.0,
        space_after=2,
    )
    table(
        s,
        x + 0.12,
        y + 0.54,
        w - 0.24,
        [
            ["Blocking gate", "Measured", "Gate"],
            ["Lake at rest — spurious velocity", "5.98e-14 m/s", "< 1e-8"],
            ["Mass conservation — volume drift", "0.000000 %", "< 0.1 %"],
            ["Ritter dam-break — RMSE vs exact", "0.0317 m", "< 0.10 m"],
        ],
        [3.74, 1.60, 1.34],
        row_h=0.228,
    )
    tb2 = textbox(s, x + 0.12, y + 1.50, w - 0.24, 0.22)
    write(
        tb2.text_frame,
        ["Both engines against the exact Ritter solution, same grid, same instant:"],
        size=9.0,
        space_after=2,
    )
    table(
        s,
        x + 0.12,
        y + 1.74,
        w - 0.24,
        [
            ["Engine", "RMSE vs exact", "Depth at dam"],
            ["JalRaksha 2D SWE", "0.0317 m", "4.532 m"],
            ["Delft3D FM (dflowfm-cli)", "0.0349 m", "4.515 m"],
        ],
        [3.74, 1.60, 1.34],
        head_fill=TEAL,
        row_h=0.228,
    )

    # ---- top right: the proof shot ---------------------------------------
    rx, rw = 7.45, 5.35
    label(s, rx, 1.32, rw, "The gates, run from the dashboard")
    _, ph = picture(s, IMG / "validation.png", rx, 1.66, rw)
    caption(
        s,
        rx,
        1.66 + ph + 0.05,
        rw,
        "All three PASS. h₀ = 10 m, t = 40 s, Δx = 10 m, frictionless flat bed; exact 4h₀/9 "
        "= 4.444 m, and the two engines agree to 0.0294 m RMSE.",
    )

    # ---- bottom: risks and strategies ------------------------------------
    by = 4.54
    lw = 6.28
    y = label(s, x, by, lw, "Potential challenges and risks", color=HAZARD)
    box(s, x, y, lw, 1.94, fill=WARM, line=WARM_EDGE)
    tb = textbox(s, x + 0.12, y + 0.10, lw - 0.24, 1.74)
    write(
        tb.text_frame,
        [
            "**Wet–dry front instability** — negative depths and spurious velocities. This "
            "fails visibly in a live demo.",
            "**30 m GLO-30 is a surface model** with no reservoir bathymetry; it flattens "
            "water and puts canopy in the channel. 1–2 cells span a Himalayan gorge.",
            "**Breach regressions disagree by 3–4×** — the dominant uncertainty term, and a "
            "documented state of the art rather than a defect in this implementation.",
            "**SPH does not scale** — uniform resolution is an acknowledged limitation "
            "(SPHERIC Grand Challenge 3); it cannot span a basin.",
            "**Debris, not clear water** — Chamoli was ~80% rock by volume, and clear-water "
            "physics understates such an event.",
            "**Manning's n is uncertain**, and inundation extent is sensitive to it.",
        ],
        size=8.5,
        bullet="▪",
        space_after=3.4,
    )

    sx, sw = 6.78, 6.17
    y = label(s, sx, by, sw, "Strategies for overcoming these challenges", color=TEAL)
    box(s, sx, y, sw, 1.94, fill=TINT, line=EDGE)
    tb = textbox(s, sx + 0.12, y + 0.10, sw - 0.24, 1.74)
    write(
        tb.text_frame,
        [
            "**Positivity-preserving scheme** — Audusse reconstruction under a CFL ≤ 1/2 "
            "bound; lake-at-rest and mass conservation gate every merge, and both pass now.",
            "**Lead with arrival time and extent**, not point depth. Every depth figure here "
            "is labelled indicative; the DEM limit is stated, not buried.",
            "**Quote the band, never the number** — Monte-Carlo over all four regressions "
            "returns p05 / median / p95 for outflow, depth and arrival time.",
            "**Confine SPH to a near-field box** (~600 m, 15 s) and hand off one-way; the far "
            "field is left to the depth-averaged solver, where it is valid.",
            "**Positioned as a Tier-1 screening layer** under CWC CDSO_GUD_DS_05: it triages "
            "which dams warrant a full Tier-2/3 study, and does not replace one.",
        ],
        size=8.5,
        bullet="▪",
        space_after=3.4,
    )


def slide5(s):
    drop_shape(s, "TextBox 8")
    stamp_team(s)

    # ---- the numbers, from one real Tehri run -----------------------------
    tiles = [
        ("1,180,102", "", ["people in the modelled domain", "(GHSL P2023A, epoch 2020)"], NAVY),
        ("1,760", "", ["reach 0.1 m depth or more — the", "population actually at risk"], HAZARD),
        ("193", "", ["have under 15 min of warning at a", "30 min lead-time assumption"], HAZARD),
        ("50.2", "min", ["for the wave to reach Koteshwar,", "13.0 km downstream"], TEAL),
        ("46.44", "km²", ["inundated envelope, severity", "index 0.94 over water"], TEAL),
    ]
    gap = 0.13
    tw = (12.60 - gap * (len(tiles) - 1)) / len(tiles)
    for i, (v, u, c, acc) in enumerate(tiles):
        stat(s, 0.35 + i * (tw + gap), 1.26, tw, 0.98, v, u, c, accent=acc, vsize=21)
    caption(
        s,
        0.35,
        2.28,
        12.60,
        "One real run, not an illustration: Tehri Dam (260 m, 3,540 MCM), central breach, "
        "600×600 cells at 200 m over a 60 km radius, 3 h simulated. Peak breach outflow "
        "229,952 m³/s (p05–p95 117,061–353,612); breach formation 29.3 min; maximum depth "
        "anywhere 99.9 m (75.1–120.6 m). Depth on a 30 m DEM is indicative — lead with arrival time.   "
        "**Target —** at ~1 min per dam measured, screening all ~6,000 NRLD large dams at "
        "Tier-1 is under a week of compute; a state's worth in an afternoon.",
    )

    # ---- left: who it reaches, then what they get -------------------------
    x, w = 0.35, 6.30
    y = label(s, x, 2.60, w, "Potential impact on the target audience", size=13.5)
    box(s, x, y, w, 1.64, fill=WHITE, line=NAVY, lw=1.4)
    tb = textbox(s, x + 0.14, y + 0.10, w - 0.28, 1.44)
    write(
        tb.text_frame,
        [
            "**NDMA / SDMA / DDMA** — depth, extent and arrival time are exactly what an "
            "evacuation order needs; .kml drops into the GIS they already run.",
            "**CWC / NDSA** — a Tier-1 screening layer under CDSO_GUD_DS_05, to triage which "
            "of India's large dams warrant a full expert study first.",
            "**NTRO / HADR planners** — scenarios for dams and landslide lakes that cannot be "
            "ground-surveyed, including reaches with no site access.",
            "**Army / BRO / NDRF** — pre-positioning, and deliberate-breach planning of the "
            "kind executed at Phuktal in 2015.",
            "**Downstream communities** — the people counted in the 1,760 above, who at "
            "Chamoli in 2021 received no warning at all.",
        ],
        size=8.4,
        bullet="▪",
        space_after=3.2,
    )

    by = label(s, x, y + 1.70, w, "Benefits of the solution (social, economic, environmental)", size=13.5)
    cw = (w - 0.24) / 3
    for i, (head, accent, body) in enumerate(
        [
            (
                "SOCIAL",
                HAZARD,
                "Arrival time is lead time, and lead time is evacuation. The exposed population "
                "is split by how many minutes each group actually has — 193 of them under 15. "
                "Works for terrain that cannot be surveyed or reached at all.",
            ),
            (
                "ECONOMIC",
                NAVY,
                "Zero licence cost, and automation replaces weeks of specialist setup per site. "
                ".shp / .kml at the 1:50,000 and 1:10,000 scales CWC's Emergency Action Plan "
                "guidelines specify — usable, not merely viewable.",
            ),
            (
                "ENVIRONMENTAL",
                TEAL,
                "Identifies cropland and ecosystems in the flood path, and supports scheduled "
                "controlled release over emergency release. Open and reproducible: anyone can "
                "re-run the cross-check and get these numbers.",
            ),
        ]
    ):
        cx = x + i * (cw + 0.12)
        box(s, cx, by, cw, 1.80, fill=TINT, line=EDGE)
        hb = textbox(s, cx + 0.10, by + 0.08, cw - 0.20, 0.24, anchor=MSO_ANCHOR.MIDDLE)
        write(
            hb.text_frame,
            [head],
            size=10.2,
            color=accent,
            align=PP_ALIGN.CENTER,
            space_after=0,
            bold=True,
        )
        bb = textbox(s, cx + 0.12, by + 0.38, cw - 0.24, 1.34)
        write(bb.text_frame, [body], size=8.2, space_after=0, line_spacing=0.97)

    # ---- right: the terrain, then where this goes next --------------------
    rx, rw = 6.95, 6.00
    label(s, rx, 2.60, rw, "Real terrain, real reach")
    iw = 3.90
    _, ph = picture(s, IMG / "terrain3d.png", rx + (rw - iw) / 2, 2.94, iw)
    caption(
        s,
        rx,
        2.96 + ph + 0.04,
        rw,
        "Bhagirathi valley below Tehri on Copernicus GLO-30, rendered in ParaView with the "
        "computed depth draped over it. Rishikesh and Haridwar report no arrival — Rishikesh "
        "sits 18 m above the nearest channel — reported as a null, not dropped.",
        size=7.2,
    )

    ry = label(s, rx, 5.44, rw, "Roadmap — done, now, next", size=11.5)
    line_h(s, rx + 0.30, ry + 0.10, rw - 0.60, color=EDGE, t=0.02)
    colw = rw / 3
    for i, (phase, accent, when, items) in enumerate(
        [
            (
                "DONE",
                TEAL,
                "Jan – Aug 2026",
                ["Solver + blocking gates", "Delft3D FM cross-check", "Dashboard, exports, SPH"],
            ),
            (
                "NOW",
                NAVY,
                "Sep 2026",
                ["Tehri + Khadakwasla runs", "Sentinel-1 validation loop", "SIH prototype review"],
            ),
            (
                "NEXT",
                HAZARD,
                "Oct 2026 →",
                ["All CWC NRLD dams", "EAP-format map packs", "SDMA field pilot"],
            ),
        ]
    ):
        cx = rx + i * colw
        node = box(s, cx + 0.16, ry + 0.02, 0.16, 0.16, fill=accent, shape=MSO_SHAPE.OVAL)
        node  # the dot on the timeline
        hb = textbox(s, cx + 0.38, ry + 0.01, colw - 0.44, 0.18, anchor=MSO_ANCHOR.MIDDLE)
        write(hb.text_frame, [phase], size=8.6, color=accent, space_after=0, bold=True)
        wb = textbox(s, cx + 0.16, ry + 0.21, colw - 0.24, 0.16)
        write(wb.text_frame, [when], size=7.0, color=MUTED, space_after=0)
        ib = textbox(s, cx + 0.16, ry + 0.40, colw - 0.24, 0.62)
        write(ib.text_frame, items, size=7.4, color=INK, bullet="▪", space_after=1.2)


def slide6(s):
    drop_shape(s, "TextBox 8")
    stamp_team(s)
    label(s, 0.35, 1.30, 12.6, "Details / Links of the reference and research work")

    top, h = 1.66, 4.36
    cols = [
        (
            0.35,
            4.14,
            "Problem domain — Indian events and regulation",
            NAVY,
            TINT,
            EDGE,
            [
                "CWC / CDSO, **Guidelines for Mapping Flood Risks Associated with Dams**, "
                "CDSO_GUD_DS_05 v1.0, 2018 — the Tier 1–3 framework this tool targets",
                "CWC, **Emergency Action Plan guidelines** — index maps 1:50,000, detailed "
                "maps 1:10,000",
                "Shugar et al. 2021, **Science** 373:300–306, doi 10.1126/science.abh4455 — "
                "Chamoli; travel times matched to under 5%",
                "Rana et al. 2023, **Natural Hazards**, doi 10.1007/s11069-023-05972-5 — an "
                "independent reconstruction of the same event",
                "Sinha 2009, **Current Science** 97:429–433 — the 2008 Kosi avulsion",
                "Costa & Schuster 1988, **GSA Bulletin** — formation and failure of natural dams",
                "CWC **NRLD 2019** and **NRSD 2025** national dam registers",
            ],
        ),
        (
            4.62,
            4.14,
            "Numerical methods",
            NAVY,
            TINT,
            EDGE,
            [
                "Toro 2001, **Shock-Capturing Methods for Free-Surface Shallow Flows**, Wiley "
                "— HLLC wave speeds, two-rarefaction estimate",
                "Audusse et al. 2004, **SIAM J. Sci. Comput.** 25:2050 — hydrostatic "
                "reconstruction, the well-balanced property this solver rests on",
                "Froehlich 2008, **J. Hydraul. Eng.** 134(12) — embankment breach regressions",
                "Xu & Zhang 2009, **J. Geotech. Geoenviron. Eng.** 135(12)",
                "Wahl 2004, **J. Hydraul. Eng.** 130(5) — uncertainty of breach parameters",
                "Maranzoni & Tomirotti 2023, **Water** 15(17):3130, doi 10.3390/w15173130 — "
                "recommends exactly this coupled 2D–3D decomposition",
                "Vacondio et al. 2020, **Comp. Particle Mech.** 8:575 — SPH Grand Challenge 3",
                "Ramachandran et al. 2021, **ACM TOMS** 47(4), doi 10.1145/3460773 — PySPH",
            ],
        ),
        (
            8.90,
            4.05,
            "Open data and software — free, none login-gated",
            TEAL,
            WARM,
            WARM_EDGE,
            [
                "**Copernicus DEM GLO-30** — 30 m, public AWS Cloud-Optimized GeoTIFFs",
                "**Sentinel-1 GRD** and **JRC Global Surface Water** via Google Earth Engine",
                "**GHSL P2023A** gridded population, epoch 2020 (JRC)",
                "**ESA WorldCover** — land cover driving Manning's n",
                "**Google Open Buildings** — CC BY 4.0",
                "**Delft3D FM**, Deltares — Suite 2026.01 HM, dflowfm-cli",
                "**PySPH** — BSD licence, developed at IIT Bombay",
                "*Deliberately excluded:* FABDEM (CC BY-NC-SA), MERIT (CC BY-NC / ODbL) and "
                "OSM (ODbL share-alike) are kept out of redistributed outputs on licence "
                "grounds.",
            ],
        ),
    ]
    for x, w, head, accent, fill, edge, items in cols:
        y = label(s, x, top, w, head, size=10.2, color=accent)
        box(s, x, y, w, h - (y - top), fill=fill, line=edge)
        tb = textbox(s, x + 0.12, y + 0.11, w - 0.24, h - (y - top) - 0.22)
        write(tb.text_frame, items, size=9.3, bullet="▪", space_after=9.5, line_spacing=0.96)

    box(s, 0.35, 6.20, 12.6, 0.56, fill=NAVY)
    tb = textbox(s, 0.52, 6.30, 12.26, 0.40, anchor=MSO_ANCHOR.MIDDLE)
    write(
        tb.text_frame,
        [
            "**Reproduce every number on this deck:**   "
            "python scripts/validate_against_delft3d.py --case ritter   |   "
            "jalraksha run --dam tehri   |   dashboard at localhost:3000",
        ],
        size=9.2,
        color=WHITE,
        space_after=0,
    )


# --------------------------------------------------------------------------- #


def main() -> None:
    prs = Presentation(str(TEMPLATE))

    # The template's own instructions cap the deck at six slides including the
    # title, so its trailing IMPORTANT INSTRUCTIONS slide has to go.
    delete_slide(prs, 6)

    for fn, s in zip((slide1, slide2, slide3, slide4, slide5, slide6), list(prs.slides)):
        fn(s)

    prs.save(str(OUTPUT))
    n = len(prs.slides._sldIdLst)
    print(f"wrote {OUTPUT.name}  ({OUTPUT.stat().st_size / 1024:.0f} KB, {n} slides)")


if __name__ == "__main__":
    main()
