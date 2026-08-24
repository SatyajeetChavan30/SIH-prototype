"""
Overflow check for the generated deck.

    python check_ppt.py

There is no LibreOffice in this environment, so the deck cannot be rendered to
images for visual inspection. Instead this measures every text frame with the
real font metrics (PIL + the actual Windows TTFs), re-does PowerPoint's
word wrap, and reports any frame whose text is taller than its shape.

Anything listed as OVERFLOW will be visibly clipped or spill out of its card in
the PDF export. Fix by shortening the text in build_ppt.py or dropping the font
size, then re-run both scripts.
"""

from __future__ import annotations

from pathlib import Path

from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
DECK = HERE / "JalRaksha_SIH2026_PS26161.pptx"

FONT_DIR = Path("C:/Windows/Fonts")
FILES = {
    ("Calibri", False, False): "calibri.ttf",
    ("Calibri", True, False): "calibrib.ttf",
    ("Calibri", False, True): "calibrii.ttf",
    ("Calibri", True, True): "calibriz.ttf",
    ("Arial", False, False): "arial.ttf",
    ("Arial", True, False): "arialbd.ttf",
    ("Times New Roman", False, False): "times.ttf",
    ("Times New Roman", True, False): "timesbd.ttf",
    ("Garamond", False, False): "GARA.TTF",
    ("Garamond", True, False): "GARABD.TTF",
}
PX_PER_PT = 8  # render at 8x for sub-point measurement accuracy
_cache: dict = {}


def font(name: str, size_pt: float, bold: bool, italic: bool):
    name = name or "Calibri"
    key = (name, bold, italic, round(size_pt, 1))
    if key in _cache:
        return _cache[key]
    fn = (
        FILES.get((name, bold, italic))
        or FILES.get((name, bold, False))
        or FILES.get((name, False, False))
        or "calibri.ttf"
    )
    f = ImageFont.truetype(str(FONT_DIR / fn), int(round(size_pt * PX_PER_PT)))
    _cache[key] = f
    return f


def measure(text: str, f) -> float:
    """Width in points."""
    return f.getlength(text) / PX_PER_PT


def wrap_lines(runs, width_pt: float) -> int:
    """Count wrapped lines for one paragraph, given [(text, font)] runs."""
    if not runs:
        return 1
    # Flatten into (word, font) tokens, preserving explicit newlines.
    tokens: list = []
    for text, f in runs:
        for i, seg in enumerate(text.split("\n")):
            if i:
                tokens.append(("\n", f))
            for j, word in enumerate(seg.split(" ")):
                if j:
                    tokens.append((" ", f))
                if word:
                    tokens.append((word, f))
    lines, cur = 1, 0.0
    pending_space = 0.0
    for tok, f in tokens:
        if tok == "\n":
            lines += 1
            cur, pending_space = 0.0, 0.0
            continue
        if tok == " ":
            pending_space = measure(" ", f)
            continue
        w = measure(tok, f)
        if cur and cur + pending_space + w > width_pt:
            lines += 1
            cur = w
        else:
            cur += pending_space + w
        pending_space = 0.0
    return lines


def frame_height_pt(tf, width_pt: float) -> float:
    """Estimated rendered text height in points."""
    total = 0.0
    for para in tf.paragraphs:
        runs = [
            (
                r.text,
                font(
                    r.font.name,
                    (r.font.size or Emu(127000)).pt,
                    bool(r.font.bold),
                    bool(r.font.italic),
                ),
            )
            for r in para.runs
            if r.text
        ]
        sizes = [(r.font.size or Emu(127000)).pt for r in para.runs if r.text]
        max_pt = max(sizes) if sizes else 12.0
        n = wrap_lines(runs, width_pt)
        spacing = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
        # PowerPoint line height ~= 1.2 * point size, scaled by line_spacing.
        total += n * max_pt * 1.2 * spacing
        total += para.space_after.pt if para.space_after is not None else 0.0
    return total


def main() -> None:
    prs = Presentation(DECK)
    problems = 0
    for si, slide in enumerate(prs.slides, 1):
        rows = []
        for sh in slide.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            tf = sh.text_frame
            avail_w = Emu(sh.width).pt - (
                (tf.margin_left.pt if tf.margin_left is not None else 7.2)
                + (tf.margin_right.pt if tf.margin_right is not None else 7.2)
            )
            avail_h = Emu(sh.height).pt - (
                (tf.margin_top.pt if tf.margin_top is not None else 3.6)
                + (tf.margin_bottom.pt if tf.margin_bottom is not None else 3.6)
            )
            need = frame_height_pt(tf, avail_w)
            ratio = need / avail_h if avail_h else 99
            label = "OVERFLOW" if ratio > 1.0 else ("tight" if ratio > 0.92 else "ok")
            if label != "ok":
                problems += 1
            rows.append((label, ratio, sh.name, need, avail_h, tf.text[:58].replace("\n", " / ")))

        rows.sort(key=lambda r: -r[1])
        print(f"\n=== SLIDE {si} " + "=" * 52)
        for label, ratio, name, need, avail, preview in rows:
            flag = {"OVERFLOW": "!!", "tight": " ~", "ok": "  "}[label]
            print(f"{flag} {ratio:5.2f}  {name:<28} {need:6.1f}/{avail:6.1f}pt  {preview}")

    # geometry sanity: anything off-canvas?
    print("\n=== OFF-CANVAS / BOUNDS " + "=" * 36)
    W, H = Emu(prs.slide_width).inches, Emu(prs.slide_height).inches
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.left is None:
                continue
            l, t = Emu(sh.left).inches, Emu(sh.top).inches
            r, b = l + Emu(sh.width).inches, t + Emu(sh.height).inches
            if l < -0.05 or t < -0.15 or r > W + 0.05 or b > H + 0.05:
                print(
                    f"   slide {si}  {sh.name:<28} "
                    f"[{l:6.2f},{t:6.2f}] -> [{r:6.2f},{b:6.2f}]  "
                    f"(canvas {W:.2f} x {H:.2f})"
                )

    print(f"\n{problems} frame(s) tight or overflowing.")


if __name__ == "__main__":
    main()
