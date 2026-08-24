"""
Build the SIH 2026 idea-submission deck for PS 26161 (JalRaksha) on top of the
official SIH template.

    python build_ppt.py

Reads : SIH2026-IDEA-Presentation-Format.pptx   (official template, unmodified)
Writes: JalRaksha_SIH2026_PS26161.pptx

Hard constraints enforced by this script, taken from the template's own
"IMPORTANT INSTRUCTIONS" slide:
  * Maximum 6 slides INCLUDING the title slide -> the instructions slide is deleted.
  * The four section titles are left exactly as the template names them
    (TECHNICAL APPROACH / FEASIBILITY AND VIABILITY / IMPACT AND BENEFITS /
    RESEARCH AND REFERENCES). Only slide 2's title is replaced, with the idea
    title, which is what that slide is for.
  * Points and diagrams, not paragraphs.
  * Must be exported to PDF before uploading to the portal.

Every factual claim below traces to docs/RESEARCH-FINDINGS.md. Edit CONTENT;
leave the layout code alone.  Inline **bold** and *italic* markup is supported
in every string.
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
TEMPLATE = HERE / "SIH2026-IDEA-Presentation-Format.pptx"
OUTPUT = HERE / "JalRaksha_SIH2026_PS26161.pptx"

# --------------------------------------------------------------------------- #
# Palette — sampled from the template so the deck matches rather than clashes.
# --------------------------------------------------------------------------- #
SIH_BLUE = RGBColor(0x00, 0x70, 0xC0)  # template footer bar
NAVY = RGBColor(0x14, 0x30, 0x5A)
INK = RGBColor(0x26, 0x2B, 0x33)
MUTED = RGBColor(0x5A, 0x63, 0x70)
TINT = RGBColor(0xED, 0xF4, 0xFB)  # card fill
TINT_EDGE = RGBColor(0xBB, 0xD6, 0xEE)  # card border
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xB5, 0x6A, 0x00)  # risk / caution
TEAL = RGBColor(0x0B, 0x74, 0x6B)  # mitigation / positive
CRIMSON = RGBColor(0xA5, 0x1E, 0x2D)  # human cost

BODY_FONT = "Calibri"

# --------------------------------------------------------------------------- #
# CONTENT
# --------------------------------------------------------------------------- #

IDEA_NAME = "JalRaksha"
TEAM_NAME = "TEAM NAME"  # <-- replace with your registered team name

TITLE_FIELDS = [
    ("Problem Statement ID", "26161"),
    (
        "Problem Statement Title",
        "Dam Break Inundation Modelling Using Hydrodynamic Modelling of any River",
    ),
    ("Theme", "Disaster Management"),
    ("PS Category", "Software"),
    ("Team ID", "[from SIH portal]"),
    ("Team Name", "[as registered on portal]"),
]

TAGLINE = (
    "Automated dam-break AND landslide-dam inundation modelling:  "
    "Smoothed Particle Hydrodynamics  x  Delft3D-class shallow water  "
    "—  quantitatively compared, satellite-validated, zero ground survey"
)

CONTENT = {
    # ---------------------------------------------------------------- slide 2
    "idea": [
        (
            "Detailed explanation of the proposed solution",
            SIH_BLUE,
            [
                "Operator enters a **dam** or a **suspected river blockage** — lat/lon, "
                "height, gross storage, breach mode. Nothing else.",
                "Pipeline auto-fetches **Copernicus GLO-30 DEM** (public AWS COGs, no "
                "login), **CWC NRLD/NRSD** dam registers, **Sentinel-1** SAR, GHSL and "
                "Open Buildings exposure.",
                "**Breach hydrograph** from Froehlich (2008) / Xu–Zhang (2009) "
                "regressions, coupled to weir outflow and reservoir drawdown.",
                "**Near field** — Smoothed Particle Hydrodynamics resolves the violent, "
                "non-hydrostatic breach jet, where depth-averaging is simply invalid.",
                "**Far field** — Delft3D-class depth-averaged shallow-water finite "
                "volume routes the wave tens of km down the valley.",
                "Outputs depth, velocity and **arrival-time** rasters to a dashboard, "
                "plus GeoTIFF / **.SHP** / **.KML**.",
            ],
        ),
        (
            "How it addresses the problem",
            NAVY,
            [
                "**(i)** Both scenario classes — engineered dam break **and** natural "
                "river blockage (Rishi Ganga, Phuktal class) — with loss-and-damage output.",
                "**(ii)** Pluggable datasets: any DEM, any hydrology source, any AOI.",
                "**(iii)** Dashboard GUI for model input and tiled output visualisation; "
                "**.shp / .kml** export as specified.",
                "**(iv)** **Google Earth Engine** Sentinel-1 SAR module for near-real-time "
                "observed flood extent.",
                "**(v)** Runs live on open Indian data — **Tehri** (260 m, 3,540 MCM, "
                "Koteshwar 13 km downstream) as the engineered case; **Chamoli 2021** as "
                "the validated natural case.",
                "**Automatic by default** — coordinates in, decision map out. No hydraulic "
                "modeller in the loop, no bathymetric survey, no site access.",
            ],
        ),
        (
            "Innovation and uniqueness of the solution",
            TEAL,
            [
                "**The comparison is the product.** The PS asks that the two models be "
                "compared; we make agreement a *measured* output — CSI / F1 on wet-cell "
                "extent, depth RMSE, arrival-time delta.",
                "**Physics-appropriate decomposition, not duplication.** SPH is "
                "uniform-resolution and cannot span a basin (SPHERIC Grand Challenge 3); "
                "depth-averaging cannot resolve a breach jet. Coupled 2D–3D is the "
                "direction the review literature itself recommends — we build it.",
                "**Validated against a published benchmark, not a pretty animation.** "
                "Shugar et al. reproduced Chamoli travel times to **<5%**. That is our "
                "target, and the pre/post 2 m DEMs are openly downloadable.",
                "**Zero-survey operation** — works on transboundary and inaccessible "
                "terrain that cannot be instrumented or entered.",
                "**Cross-checked against the incumbent:** the same reach also runs through "
                "HEC-RAS, the tool CWC adopted for DRIP.",
            ],
        ),
    ],
    # ---------------------------------------------------------------- slide 3
    "tech": [
        ("Compute", "Python 3.11 · NumPy · Numba / CuPy (GPU) · xarray · Dask"),
        (
            "Near-field solver",
            "**PySPH** — BSD, pip-installable, developed at IIT Bombay; ships dam-break "
            "and shallow-water SPH examples with bundled validation data. "
            "DualSPHysics (LGPL, CUDA) for GPU scale-up.",
        ),
        (
            "Far-field solver",
            "Delft3D FM (Deltares, AGPL/GPL) driven by **hydromt_delft3dfm · "
            "hydrolib-core · dfm_tools**; plus an in-house 2D shallow-water finite "
            "volume — HLLC Riemann, MUSCL, well-balanced, wetting/drying.",
        ),
        ("Cross-check", "HEC-RAS 2D — the CWC / DRIP baseline"),
        ("Geospatial", "GDAL · rasterio · geopandas · shapely · simplekml"),
        ("Earth observation", "Google Earth Engine Python API · geemap · Sentinel-1 GRD"),
        (
            "Data",
            "Copernicus GLO-30 (AWS COG) · CWC NRLD 2019 / NRSD 2025 · GHSL · "
            "Google Open Buildings · ESA WorldCover · JRC Global Surface Water",
        ),
        ("Application", "FastAPI · Celery · React + deck.gl / MapLibre GL"),
        ("Export", "GeoTIFF · **.SHP** · **.KML / KMZ**"),
    ],
    "flow": [
        (
            "1",
            "INPUT",
            "Dam or suspected blockage — lat/lon, height, gross storage, breach mode",
            SIH_BLUE,
        ),
        (
            "2",
            "AUTO-INGEST",
            "Copernicus GLO-30 COGs · CWC NRLD/NRSD · Sentinel-1 · GHSL · WorldCover",
            SIH_BLUE,
        ),
        (
            "3",
            "BREACH HYDROGRAPH",
            "Froehlich 2008 / Xu–Zhang 2009 regressions → weir outflow + reservoir drawdown",
            SIH_BLUE,
        ),
        ("SPLIT", None, None, None),
        (
            "5",
            "COMPARE   ←  the PS's core ask",
            "CSI / F1 on inundation extent · depth RMSE · arrival-time delta",
            TEAL,
        ),
        (
            "6",
            "LOSS AND DAMAGE",
            "Population, buildings, roads and bridges inside the flood path",
            NAVY,
        ),
        ("7", "OUTPUT", "Dashboard · arrival-time map · GeoTIFF / .SHP / .KML", NAVY),
        (
            "8",
            "GEE VALIDATION LOOP",
            "Sentinel-1 SAR observed extent scored against the simulation ↻",
            AMBER,
        ),
    ],
    "flow_split": [
        ("4a", "NEAR FIELD — SPH", "Lagrangian, mesh-free. The non-hydrostatic breach jet."),
        (
            "4b",
            "FAR FIELD — Delft3D-class 2D SWE",
            "Eulerian finite volume. HLLC + wetting/drying.",
        ),
    ],
    # ---------------------------------------------------------------- slide 4
    "feas": [
        (
            "Analysis of the feasibility of the idea",
            TEAL,
            [
                "**Every input verified free, open and login-free** — Copernicus GLO-30 "
                "tiles serve as range-readable COGs from public AWS; Sentinel-1, GHSL and "
                "Open Buildings are native Earth Engine assets.",
                "**Both named models are genuinely obtainable.** PySPH is BSD and "
                "`pip install`-able with dam-break cases in the box; Delft3D FM is "
                "AGPL/GPL with a maintained Python toolchain.",
                "The far-field solver is **textbook numerics** — HLLC Riemann solver, "
                "MUSCL reconstruction, well-balanced source terms. GPU-accelerated it runs "
                "a 30 m reach in minutes on a laptop.",
                "**A hard accuracy target already exists:** Shugar et al. matched Chamoli "
                "travel times to **<5%**, and the pre- and post-event 2 m DEMs are on "
                "Zenodo. Measured discharge 8,200–14,200 m³/s at 15 km.",
                "**Scope control** — SPH confined to a near-field box, so particle count "
                "stays tractable; published real-field cases run 1–4 million particles.",
                "Skills needed: Python, GIS, numerical methods. **No proprietary licence, "
                "no field survey, no bathymetric campaign.**",
            ],
        ),
        (
            "Potential challenges and risks",
            AMBER,
            [
                "**Numerical instability at the wet–dry front** — negative depths, spurious "
                "velocities, mass loss. This fails *visibly* in a live demo.",
                "**SPH does not scale.** Uniform resolution is an acknowledged limitation "
                "(SPHERIC GC3); it cannot span a basin.",
                "**No reservoir bathymetry in any DEM**, and GLO-30 is a *surface* model — "
                "it flattens water and puts canopy in the channel.",
                "**30 m cells give 1–2 cells across a Himalayan gorge**, so point depths "
                "there are fiction, and a cell straddling an embankment lowers the "
                "effective crest by metres.",
                "**Manning's roughness is uncertain** and inundation extent is sensitive to it.",
                "**Debris, not clear water.** Chamoli was ~80% rock by volume; clear-water "
                "physics understates such an event.",
                "**Delft3D kernel build risk** — Fortran/C++ toolchain, only unofficial "
                "Docker images. **Liability** — a map an officer acts on has consequences.",
            ],
        ),
        (
            "Strategies for overcoming these challenges",
            SIH_BLUE,
            [
                "Well-balanced HLLC with **hydrostatic reconstruction** (Audusse 2004), a "
                "depth-positivity limiter and adaptive CFL. **Verify against the Ritter "
                "analytical solution before touching terrain** — correctness first.",
                "**Couple, do not choose:** SPH near field hands off to shallow-water far "
                "field across a coupling boundary — the review literature's own recommendation.",
                "Reservoir volume from **NRLD gross storage** fitted to a synthetic "
                "bathymetric wedge; channel burned in from MERIT Hydro. Report a "
                "**sensitivity band, never a single line**.",
                "Report **arrival time and inundation envelope** — which 30 m data supports "
                "— and label point depths as indicative. Publish a Manning's-*n* ensemble.",
                "Debris **flagged explicitly** and approximated with bulked density and "
                "roughness. Stated as a caveat, not hidden.",
                "**Our own solver carries the live demo**; Delft3D FM is the offline "
                "reference run, so a build failure cannot break the demonstration.",
                "Every output stamped **Tier-1 screening-grade** per CWC CDSO_GUD_DS_05, "
                "with uncertainty printed on the map itself.",
            ],
        ),
    ],
    # ---------------------------------------------------------------- slide 5
    "stats": [
        ("204", "dead or missing", "Chamoli, 7 Feb 2021 — no warning reached anyone", CRIMSON),
        ("0", "deaths", "Phuktal, 7 May 2015 — 30–35 Mm³ released, ~3,000 evacuated", TEAL),
        (
            "2.3 M+",
            "people affected",
            "Kosi, 18 Aug 2008 — breached far below design discharge",
            AMBER,
        ),
    ],
    "stat_note": (
        "The variable between these three rows is not the hydrology. It is "
        "whether anyone knew where the water would go, and when."
    ),
    "impact": [
        (
            "Potential impact on the target audience",
            NAVY,
            [
                "**NTRO / HADR planners** — scenario generation for dams and natural lakes "
                "that cannot be ground-surveyed, including transboundary reaches that are, "
                "in ORF's words on the Yarlung Tsangpo, *\"presently a blind spot for "
                'India"*, where the primary risk is dam **failure**, not diversion.',
                "**NDMA / SDMA / DDMA** — depth, extent and *arrival time* are exactly what "
                "an evacuation order needs; .kml drops into the GIS they already run.",
                "**CWC / NDSA** — a Tier-1 screening layer (CDSO_GUD_DS_05) to triage which "
                "dams warrant a full expert study first.",
                "**Army / BRO / NDRF** — pre-positioning, and deliberate-breach planning of "
                "the kind the Army executed at Phuktal in May 2015.",
                "**Downstream communities** — 29 villages, ~4,000 people and ~50 bridges "
                "were at risk at Phuktal alone.",
            ],
        ),
        (
            "Benefits of the solution",
            TEAL,
            [
                "**Social** — arrival-time maps become lead time, and lead time becomes "
                "evacuation. Phuktal: ~3,000 moved, zero deaths.",
                "**Economic** — zero licence cost against commercial suites, and automation "
                "replaces weeks of specialist setup per site.",
                "**Environmental** — identifies ecosystems and cropland at risk; supports "
                "scheduled controlled release instead of emergency release.",
                "**Strategic** — a no-survey capability for terrain that cannot be "
                "physically accessed or instrumented.",
                "**Institutional** — .shp / .kml at 1:50,000 and 1:10,000 matches CWC EAP "
                "map specifications, so results are usable rather than merely viewable.",
                "**Scientific** — the SPH vs depth-averaged agreement metrics are a "
                "publishable contribution, not a repackaged UI.",
            ],
        ),
    ],
    # ---------------------------------------------------------------- slide 6
    "refs": [
        (
            "Problem domain — Indian events and regulation",
            SIH_BLUE,
            [
                "CWC / CDSO, **Guidelines for Mapping Flood Risks Associated with Dams**, "
                "CDSO_GUD_DS_05 v1.0, Jan 2018 — HEC-RAS adopted for DRIP; Tier 1–3 framework",
                "CWC, **Emergency Action Plan guidelines**, cwc.gov.in — index maps "
                "1:50,000; detailed maps 1:10,000 at 0.5 m contours",
                "Shugar et al. 2021, *Science* 373:300–306, doi 10.1126/science.abh4455 — "
                "Chamoli rock-and-ice avalanche; 26.9 Mm³; travel times matched to **<5%**",
                "Rana et al. 2023, *Natural Hazards*, doi 10.1007/s11069-023-05972-5 — "
                "independent HEC-RAS reconstruction of the same event",
                "Sinha 2009, **The Great Avulsion of Kosi on 18 August 2008**, "
                "*Current Science* 97:429–433",
                "Costa & Schuster 1988, **The formation and failure of natural dams**, "
                "*GSA Bulletin*",
                "NRSC / ISRO Cartosat-2 monitoring of the Phuktal blockage, Jan–Mar 2015",
                "CWC **NRLD 2019** and **NRSD 2025** national dam registers (cwc.gov.in)",
            ],
        ),
        (
            "Numerical methods",
            NAVY,
            [
                "Gingold & Monaghan 1977, *MNRAS* 181:375 — origin of SPH; Lucy 1977, *AJ* 82:1013",
                "Maranzoni & Tomirotti 2023, **3D Numerical Modelling of Real-Field "
                "Dam-Break Flows**, *Water* 15(17):3130, doi 10.3390/w15173130 — recommends "
                "coupled 2D–3D",
                "Vacondio et al. 2020, **Grand challenges for SPH numerical schemes**, "
                "*Comp. Particle Mech.* 8:575–588 — GC3: uniform resolution blocks "
                "multiscale use",
                "Vacondio et al. 2011, **SPH Modeling of Shallow Flow with Open Boundaries "
                "for Practical Flood Simulation**, *J. Hydraul. Eng.*",
                "Prakash, Rothauge & Cleary 2014, **Modelling the impact of dam failure "
                "scenarios on flood inundation using SPH**, *Appl. Math. Modelling*",
                "Ramachandran et al. 2021, **PySPH: A Python-based framework for SPH**, "
                "*ACM TOMS* 47(4), doi 10.1145/3460773  (IIT Bombay)",
                "Toro 2001, *Shock-Capturing Methods for Free-Surface Shallow Flows*, Wiley; "
                "Audusse et al. 2004, *SIAM J. Sci. Comput.*; Liang & Marche 2009, *AWR*",
                "Froehlich 2008 breach parameters, *J. Hydraul. Eng.*; Xu & Zhang 2009, *JGGE*",
            ],
        ),
        (
            "Software, data and Earth observation",
            TEAL,
            [
                "**Delft3D** — github.com/Deltares/Delft3D (AGPL-3.0 / GPL-3.0); "
                "oss.deltares.nl; `hydromt_delft3dfm`, `hydrolib-core`, `dfm_tools`",
                "**PySPH** — pysph.readthedocs.io (BSD); **DualSPHysics** — dual.sphysics.org "
                "(LGPL, CUDA); GPUSPH (GPLv3)",
                "**HEC-RAS** — hec.usace.army.mil, incl. official dam-break study guidance",
                "**Copernicus DEM GLO-30** — public COGs at s3://copernicus-dem-30m; "
                "GEE `COPERNICUS/DEM/GLO30_2024_1`. EGM2008 datum, <4 m vertical (90% LE)",
                "**Chamoli validation DEMs** — pre-event 2 m zenodo.org/record/4554647, "
                "post-event zenodo.org/record/4558692 (CC BY-NC 4.0)",
                "GEE assets — `COPERNICUS/S1_GRD` · `JRC/GSW1_4/GlobalSurfaceWater` · "
                "`JRC/GHSL/P2023A/GHS_BUILT_C` · `WorldPop/GP/100m/pop` · "
                "`GOOGLE/Research/open-buildings/v3/polygons` · `ESA/WorldCover/v200`",
                "**Global Dam Watch v1.0** (CC BY 4.0), doi 10.6084/m9.figshare.25988293 — "
                "41,145 barriers, HydroSHEDS-harmonised, fallback inventory",
                "**Malpasset 1959** — the standard real-world dam-break validation case; "
                "**Ritter/Stoker** analytical solution for solver verification",
            ],
        ),
    ],
}

FOOTNOTE = (
    "Sources verified Aug 2026. Two dating notes on the PS's reference list: the "
    "Phuktal blockage formed 31 Dec 2014 and was breached 7 May 2015; no Nov 2021 "
    "Indian river blockage is on record — the likely referent is the 29 Oct 2021 "
    "Kameng (Warriyang) sediment event, in which no dam formed. Chamoli was a "
    "rock-and-ice avalanche, not a GLOF."
)


# --------------------------------------------------------------------------- #
# Text helpers
# --------------------------------------------------------------------------- #

_BOLD = re.compile(r"\*\*(.+?)\*\*")
_ITAL = re.compile(r"(?<!\*)\*([^*]+?)\*(?!\*)")


def _segments(text: str):
    """Split a string into (chunk, bold, italic) triples using **bold** / *italic*."""
    out, pos = [], 0
    for m in _BOLD.finditer(text):
        if m.start() > pos:
            out.append((text[pos : m.start()], False))
        out.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], False))

    final = []
    for chunk, bold in out:
        p = 0
        for m in _ITAL.finditer(chunk):
            if m.start() > p:
                final.append((chunk[p : m.start()], bold, False))
            final.append((m.group(1), bold, True))
            p = m.end()
        if p < len(chunk):
            final.append((chunk[p:], bold, False))
    return final


def write(
    tf,
    lines,
    *,
    size=10.5,
    color=INK,
    bullet=None,
    space_after=4,
    line_spacing=0.92,
    align=PP_ALIGN.LEFT,
    font=BODY_FONT,
):
    """Fill a text frame. `lines` is a list of markup strings."""
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(space_after)
        para.line_spacing = line_spacing
        text = f"{bullet} {line}" if bullet else line
        for chunk, bold, italic in _segments(text):
            run = para.add_run()
            run.text = chunk
            run.font.size = Pt(size)
            run.font.name = font
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
    return tf


def box(slide, x, y, w, h, *, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.06):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = adj
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.9)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    return sp


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = 0
    return tb


def card(slide, x, y, w, h, heading, accent, bullets, *, size=9.6, head_size=11.5):
    """Tinted card: coloured heading strip + bulleted body."""
    box(slide, x, y, w, h, fill=TINT, line=TINT_EDGE)
    box(slide, x, y, w, 0.34, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    hd = textbox(slide, x + 0.10, y + 0.045, w - 0.20, 0.26, anchor=MSO_ANCHOR.MIDDLE)
    write(hd.text_frame, [heading], size=head_size, color=WHITE, space_after=0)
    bd = textbox(slide, x + 0.13, y + 0.43, w - 0.26, h - 0.55)
    write(bd.text_frame, bullets, size=size, bullet="▪", space_after=4.5)


def drop_shape(slide, name):
    for sh in list(slide.shapes):
        if sh.name == name:
            sh._element.getparent().remove(sh._element)
            return True
    return False


def _reset(tf):
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    tf.paragraphs[0].clear()


def set_shape_text(slide, name, lines, **kw):
    for sh in slide.shapes:
        if sh.name == name and sh.has_text_frame:
            _reset(sh.text_frame)
            write(sh.text_frame, lines, **kw)
            return sh
    return None


def stamp_team(slide):
    for sh in slide.shapes:
        if sh.name.startswith("Oval") and sh.has_text_frame:
            _reset(sh.text_frame)
            write(
                sh.text_frame,
                [TEAM_NAME],
                size=8.5,
                color=INK,
                align=PP_ALIGN.CENTER,
                space_after=0,
            )


def retitle(slide, text, size=32):
    for sh in slide.shapes:
        if sh.name.startswith("Title"):
            _reset(sh.text_frame)
            run = sh.text_frame.paragraphs[0].add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.name = "Times New Roman"
            return


def arrow(slide, cx, y, h=0.10, w=0.20, color=SIH_BLUE):
    sp = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(cx - w / 2), Inches(y), Inches(w), Inches(h)
    )
    sp.rotation = 180
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False


# --------------------------------------------------------------------------- #
# Build
# --------------------------------------------------------------------------- #


def main() -> None:
    prs = Presentation(TEMPLATE)

    # --- drop the "IMPORTANT INSTRUCTIONS" slide (6-slide cap) ---------------
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    if len(ids) == 7:
        prs.part.drop_rel(ids[6].rId)
        sldIdLst.remove(ids[6])

    s1, s2, s3, s4, s5, s6 = prs.slides
    xs = [0.32, 4.67, 9.02]  # three-column x origins
    CW = 4.15  # column width

    # ======================================================== SLIDE 1: title
    set_shape_text(
        s1,
        "TextBox 9",
        [f"**{label} —** {value}" for label, value in TITLE_FIELDS],
        size=15,
        color=NAVY,
        font="Arial",
        space_after=8,
        line_spacing=1.0,
    )
    set_shape_text(
        s1,
        "Subtitle 3",
        [IDEA_NAME.upper()],
        size=26,
        color=SIH_BLUE,
        font="Times New Roman",
        space_after=0,
    )

    # ==================================================== SLIDE 2: idea title
    retitle(s2, IDEA_NAME, size=34)
    drop_shape(s2, "TextBox 8")

    band = box(
        s2, 0.32, 1.26, 12.70, 0.62, fill=SIH_BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.30
    )
    band.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(band.text_frame, [TAGLINE], size=12.0, color=WHITE, align=PP_ALIGN.CENTER, space_after=0)

    for (heading, accent, bullets), x in zip(CONTENT["idea"], xs):
        card(s2, x, 2.04, CW, 4.78, heading, accent, bullets, size=11.4)

    # =============================================== SLIDE 3: technical approach
    drop_shape(s3, "TextBox 8")

    box(s3, 0.32, 1.32, 4.00, 5.50, fill=TINT, line=TINT_EDGE)
    box(s3, 0.32, 1.32, 4.00, 0.34, fill=SIH_BLUE, shape=MSO_SHAPE.RECTANGLE)
    hd = textbox(s3, 0.42, 1.365, 3.80, 0.26, anchor=MSO_ANCHOR.MIDDLE)
    write(hd.text_frame, ["Technologies to be used"], size=11.5, color=WHITE, space_after=0)
    tb = textbox(s3, 0.45, 1.74, 3.74, 5.00)
    write(
        tb.text_frame,
        [f"**{k}**\n{v}" for k, v in CONTENT["tech"]],
        size=9.4,
        space_after=5.0,
        line_spacing=0.94,
    )

    hd = textbox(s3, 4.62, 1.32, 8.42, 0.30)
    write(
        hd.text_frame,
        ["Methodology and process for implementation"],
        size=11.5,
        color=NAVY,
        space_after=0,
    )

    fx, fw, y = 4.62, 8.42, 1.70
    for num, label, detail, colour in CONTENT["flow"]:
        if num == "SPLIT":
            hw = (fw - 0.20) / 2
            for i, (n, lab, det) in enumerate(CONTENT["flow_split"]):
                sp = box(s3, fx + i * (hw + 0.20), y, hw, 0.68, fill=WHITE, line=SIH_BLUE)
                sp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                write(
                    sp.text_frame,
                    [f"**{n}   {lab}**", det],
                    size=9.0,
                    space_after=1,
                    line_spacing=0.9,
                )
            y += 0.68 + 0.13
            arrow(s3, fx + fw / 2, y - 0.115)
            continue

        sp = box(s3, fx, y, fw, 0.50, fill=colour)
        sp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(
            sp.text_frame,
            [f"**{num}   {label}**   —   {detail}"],
            size=9.4,
            color=WHITE,
            space_after=0,
            line_spacing=0.9,
        )
        y += 0.50 + 0.13
        if num != "8":
            arrow(s3, fx + fw / 2, y - 0.115)

    # ============================================ SLIDE 4: feasibility & viability
    drop_shape(s4, "TextBox 8")
    for (heading, accent, bullets), x in zip(CONTENT["feas"], xs):
        card(s4, x, 1.32, CW, 5.50, heading, accent, bullets, size=11.0)

    # ================================================ SLIDE 5: impact & benefits
    drop_shape(s5, "TextBox 8")

    for i, (big, mid, sub, colour) in enumerate(CONTENT["stats"]):
        sp = box(s5, xs[i], 1.30, CW, 0.90, fill=WHITE, line=colour)
        sp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = sp.text_frame
        p1 = tf.paragraphs[0]
        p1.alignment, p1.space_after, p1.line_spacing = PP_ALIGN.CENTER, Pt(0), 0.86
        for txt, sz, col in ((big + "  ", 26, colour), (mid, 11.5, INK)):
            r = p1.add_run()
            r.text = txt
            r.font.size, r.font.bold, r.font.name = Pt(sz), True, BODY_FONT
            r.font.color.rgb = col
        p2 = tf.add_paragraph()
        p2.alignment, p2.space_after = PP_ALIGN.CENTER, Pt(0)
        r = p2.add_run()
        r.text = sub
        r.font.size, r.font.name = Pt(8.2), BODY_FONT
        r.font.color.rgb = MUTED

    note = textbox(s5, 0.32, 2.28, 12.70, 0.32)
    write(
        note.text_frame,
        [CONTENT["stat_note"]],
        size=11.5,
        color=NAVY,
        align=PP_ALIGN.CENTER,
        space_after=0,
    )

    x = 0.32
    for (heading, accent, bullets), w in zip(CONTENT["impact"], (6.28, 6.42)):
        card(s5, x, 2.70, w, 4.12, heading, accent, bullets, size=12.0)
        x += w + 0.20

    # =========================================== SLIDE 6: research & references
    drop_shape(s6, "TextBox 8")
    for (heading, accent, bullets), x in zip(CONTENT["refs"], xs):
        card(s6, x, 1.32, CW, 5.10, heading, accent, bullets, size=9.4, head_size=10.5)
    fn = textbox(s6, 0.32, 6.48, 12.70, 0.40)
    write(fn.text_frame, [FOOTNOTE], size=8.0, color=MUTED, space_after=0, line_spacing=0.92)

    # --- team name on every content slide -----------------------------------
    for s in (s2, s3, s4, s5, s6):
        stamp_team(s)

    prs.save(OUTPUT)
    print(f"wrote {OUTPUT.name}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
